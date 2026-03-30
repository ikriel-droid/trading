from __future__ import annotations

import argparse
import json
import textwrap
from pathlib import Path
from typing import Iterable

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

PROJECT_ROOT = Path(__file__).resolve().parents[1]
DIST_ROOT = PROJECT_ROOT / "dist" / "small-live-validation-guide"
ASSET_ROOT = DIST_ROOT / "assets"


def _pick_font(candidates: Iterable[str], size: int) -> ImageFont.ImageFont:
    for candidate in candidates:
        path = Path(candidate)
        if path.exists():
            return ImageFont.truetype(str(path), size=size)
    return ImageFont.load_default()


TITLE_FONT = _pick_font(
    [
        r"C:\Windows\Fonts\malgunbd.ttf",
        r"C:\Windows\Fonts\malgun.ttf",
    ],
    42,
)
BODY_FONT = _pick_font(
    [
        r"C:\Windows\Fonts\malgun.ttf",
    ],
    28,
)
CODE_FONT = _pick_font(
    [
        r"C:\Windows\Fonts\consolab.ttf",
        r"C:\Windows\Fonts\consola.ttf",
        r"C:\Windows\Fonts\malgun.ttf",
    ],
    24,
)


def _load_json(path: Path) -> dict:
    with open(path, "r", encoding="utf-8-sig") as handle:
        return json.load(handle)


def _wrap_lines(lines: Iterable[str], width: int) -> list[str]:
    wrapped: list[str] = []
    for line in lines:
        pieces = textwrap.wrap(line, width=width, break_long_words=False, replace_whitespace=False)
        wrapped.extend(pieces or [""])
    return wrapped


def _draw_terminal_screenshot(path: Path, title: str, lines: list[str]) -> Path:
    width = 1600
    padding = 36
    line_height = 42
    body_lines = _wrap_lines(lines, width=74)
    height = 120 + padding * 2 + line_height * max(5, len(body_lines))
    image = Image.new("RGB", (width, height), "#0f1720")
    draw = ImageDraw.Draw(image)

    draw.rounded_rectangle((18, 18, width - 18, height - 18), radius=28, fill="#0b1220", outline="#1f2a37", width=3)
    draw.rounded_rectangle((36, 36, width - 36, 100), radius=20, fill="#152033")
    draw.ellipse((64, 58, 84, 78), fill="#fb7185")
    draw.ellipse((96, 58, 116, 78), fill="#fbbf24")
    draw.ellipse((128, 58, 148, 78), fill="#34d399")
    draw.text((180, 48), title, font=BODY_FONT, fill="#dbeafe")

    y = 132
    for line in body_lines:
        draw.text((70, y), line, font=CODE_FONT, fill="#e5eefc")
        y += line_height

    image.save(path)
    return path


def _draw_note_screenshot(path: Path, title: str, bullets: list[str], accent: str = "#00c2a8") -> Path:
    width = 1600
    padding = 42
    bullet_lines: list[str] = []
    for bullet in bullets:
        wrapped = _wrap_lines([bullet], width=60)
        bullet_lines.extend(wrapped)
    height = 180 + padding * 2 + 42 * max(4, len(bullet_lines))
    image = Image.new("RGB", (width, height), "#f5fbfb")
    draw = ImageDraw.Draw(image)

    draw.rounded_rectangle((18, 18, width - 18, height - 18), radius=28, fill="#ffffff", outline="#d2ece8", width=3)
    draw.rounded_rectangle((48, 48, width - 48, 122), radius=18, fill=accent)
    draw.text((82, 66), title, font=TITLE_FONT, fill="#ffffff")

    y = 170
    for bullet in bullets:
        wrapped = _wrap_lines([bullet], width=58)
        draw.text((90, y), "-", font=BODY_FONT, fill="#0f3d3b")
        draw.text((128, y), wrapped[0], font=BODY_FONT, fill="#173533")
        y += 42
        for extra in wrapped[1:]:
            draw.text((128, y), extra, font=BODY_FONT, fill="#173533")
            y += 42

    image.save(path)
    return path


def _add_title(slide, title: str, subtitle: str) -> None:
    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.5), Inches(8.7), Inches(1.0))
    title_frame = title_box.text_frame
    title_frame.clear()
    para = title_frame.paragraphs[0]
    para.text = title
    para.font.size = Pt(28)
    para.font.bold = True
    para.font.color.rgb = RGBColor(8, 71, 69)

    subtitle_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.22), Inches(11.2), Inches(0.6))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.clear()
    para = subtitle_frame.paragraphs[0]
    para.text = subtitle
    para.font.size = Pt(12)
    para.font.color.rgb = RGBColor(84, 109, 111)


def _add_bullets(slide, bullets: list[str], left: float, top: float, width: float, height: float) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.word_wrap = True
    frame.clear()
    for index, bullet in enumerate(bullets):
        para = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        para.text = bullet
        para.font.size = Pt(18)
        para.font.color.rgb = RGBColor(24, 56, 54)
        para.level = 0
        para.space_after = Pt(8)


def _add_image(slide, image_path: Path, left: float, top: float, width: float) -> None:
    slide.shapes.add_picture(str(image_path), Inches(left), Inches(top), width=Inches(width))


def _add_badge(slide, text: str, left: float, top: float, width: float, fill_rgb: tuple[int, int, int]) -> None:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(0.45))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*fill_rgb)
    shape.line.color.rgb = RGBColor(*fill_rgb)
    tf = shape.text_frame
    tf.clear()
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    para.text = text
    para.font.size = Pt(11)
    para.font.bold = True
    para.font.color.rgb = RGBColor(255, 255, 255)


def build_presentation(output_path: Path, readiness_path: Path) -> dict[str, str]:
    readiness = _load_json(readiness_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    ASSET_ROOT.mkdir(parents=True, exist_ok=True)

    blockers = readiness.get("blockers") or []
    next_steps = readiness.get("next_manual_steps") or []

    env_image = _draw_terminal_screenshot(
        ASSET_ROOT / "step1-env.png",
        ".env 파일 예시",
        [
            "UPBIT_ACCESS_KEY=실제업비트AccessKey",
            "UPBIT_SECRET_KEY=실제업비트SecretKey",
            "DISCORD_WEBHOOK_URL=",
            "",
            "# 파일 위치",
            r"C:\Users\Administrator\.vscode\cli\upbit-auto-trader\.env",
        ],
    )
    config_image = _draw_terminal_screenshot(
        ASSET_ROOT / "step2-config.png",
        "config.live.micro.json 핵심만 바꾸기",
        [
            '{',
            '  "market": "KRW-BTC",',
            '  "upbit": {',
            '    "market": "KRW-BTC",',
            '    "access_key": "${UPBIT_ACCESS_KEY}",',
            '    "secret_key": "${UPBIT_SECRET_KEY}",',
            '    "live_enabled": true',
            "  }",
            "}",
        ],
    )
    bootstrap_image = _draw_terminal_screenshot(
        ASSET_ROOT / "step3-bootstrap.png",
        "한 번에 준비하기",
        [
            r".\bootstrap_small_live_validation.cmd -ConfigPath config.live.micro.json -Market KRW-BTC",
            "",
            "# 이 명령이 자동으로 하는 일",
            "1) 최근 캔들 CSV 저장",
            "2) data/live-state.json 생성 또는 복구",
            "3) readiness 스냅샷 재생성",
            "4) support bundle 재생성",
        ],
    )
    live_run_image = _draw_terminal_screenshot(
        ASSET_ROOT / "step4-live-run.png",
        "실제 소액 실거래는 터미널 2개",
        [
            "# 터미널 A",
            r".\.venv\Scripts\python.exe -m upbit_auto_trader.main listen-private --config config.live.micro.json --state data\live-state.json --market KRW-BTC --max-events 20",
            "",
            "# 터미널 B",
            r".\.venv\Scripts\python.exe -m upbit_auto_trader.main run-live-daemon --config config.live.micro.json --state data\live-state.json --warmup-csv data\live_krw_btc_15m.csv --max-loops 20 --reconcile-every-loops 1",
        ],
    )
    evidence_image = _draw_terminal_screenshot(
        ASSET_ROOT / "step5-evidence.png",
        "마지막 증빙 3개 저장",
        [
            r".\.venv\Scripts\python.exe -m upbit_auto_trader.main session-report --config config.live.micro.json --state data\live-state.json --mode live --label live-micro-validation --keep-latest 20",
            r".\build_control_room_support_bundle.cmd -StatePath data\live-state.json -CreateZip -ZipPath dist\upbit-control-room-support-live-validation.zip",
            r".\.venv\Scripts\python.exe -m upbit_auto_trader.main release-status --config config.live.micro.json",
        ],
    )
    blocker_image = _draw_note_screenshot(
        ASSET_ROOT / "blockers.png",
        "지금 막혀 있는 이유",
        [
            "현재 readiness 기준 blocker: {0}".format(", ".join(blockers) if blockers else "없음"),
            "즉, 지금은 아직 실주문 단계로 바로 가면 안 됨",
            "핵심은 .env, live_enabled=true, live-state.json 3개를 먼저 맞추는 것",
        ],
        accent="#0d9488",
    )
    next_steps_image = _draw_note_screenshot(
        ASSET_ROOT / "next-steps.png",
        "이 PPT대로 하면 되는 순서",
        [
            "1. .env 만들기",
            "2. config.live.micro.json 만들기",
            "3. bootstrap_small_live_validation.cmd 실행",
            "4. readiness.json에서 blocker 사라졌는지 확인",
            "5. 아주 작은 금액으로 live 검증",
        ],
        accent="#14b8a6",
    )

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    def add_bg(slide) -> None:
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(245, 251, 251)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    _add_title(slide, "업비트 소액 실거래 검증", "진짜로 해야 하는 것만 5단계로 줄인 설명")
    _add_badge(slide, "전체 완성 99%", 10.9, 0.45, 1.4, (13, 148, 136))
    _add_bullets(
        slide,
        [
            "이 PPT는 `Small live validation completed` 체크를 위해 필요한 최소 경로만 정리한 버전입니다.",
            "실거래 주문은 실제 돈이 나가므로, 키만 대신 넣는 수준까지는 도와줄 수 있어도 주문 클릭은 사람이 직접 확인하고 해야 합니다.",
            "복잡한 준비 단계는 `bootstrap_small_live_validation.cmd`로 한 번에 줄였습니다.",
        ],
        0.75,
        1.7,
        5.4,
        2.3,
    )
    _add_image(slide, blocker_image, 6.25, 1.55, 6.2)
    _add_image(slide, next_steps_image, 0.75, 4.25, 11.7)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    _add_title(slide, "1. .env 파일 만들기", "키는 여기 한 곳에만 넣으면 됩니다")
    _add_bullets(
        slide,
        [
            "프로젝트 폴더에서 `.env.example`을 복사해서 `.env`를 만듭니다.",
            "Upbit Access Key / Secret Key를 그대로 넣습니다.",
            "이 단계가 끝나야 `access_key_missing`, `secret_key_missing` blocker가 사라집니다.",
        ],
        0.75,
        1.7,
        4.8,
        2.0,
    )
    _add_image(slide, env_image, 5.3, 1.55, 7.0)
    _add_badge(slide, "파일 위치: 프로젝트 루트", 0.78, 4.1, 2.3, (15, 118, 110))
    _add_bullets(
        slide,
        [
            r"복사 명령: Copy-Item .env.example .env",
            r"편집 명령: notepad .env",
        ],
        0.78,
        4.7,
        4.1,
        1.3,
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    _add_title(slide, "2. live 설정 파일 만들기", "예제 설정을 복사해서 실거래 전용 파일로 따로 씁니다")
    _add_bullets(
        slide,
        [
            "기본 파일은 `config.example.json`이고, 실거래용은 `config.live.micro.json`처럼 따로 만드는 게 안전합니다.",
            "`market`과 `upbit.market`은 같은 값으로 맞춥니다.",
            "`upbit.live_enabled`는 실제 검증 직전에만 `true`로 바꿉니다.",
        ],
        0.75,
        1.7,
        4.9,
        2.2,
    )
    _add_image(slide, config_image, 5.4, 1.5, 6.9)
    _add_bullets(
        slide,
        [
            "중요: 테스트할 코인 잔고가 이미 있으면 live bootstrap이 막힐 수 있습니다.",
            "예: `KRW-BTC`로 할 거면 BTC 잔고가 0인 상태가 가장 깔끔합니다.",
        ],
        0.78,
        4.45,
        4.4,
        1.4,
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    _add_title(slide, "3. 복잡한 준비는 한 번에", "캔들 저장 + live-state 생성 + readiness 재점검")
    _add_bullets(
        slide,
        [
            "이제는 긴 inline Python을 칠 필요 없이 아래 명령 하나만 쓰면 됩니다.",
            "성공하면 `data/live-state.json`과 `dist/live-validation/small-live-validation-readiness.json`이 같이 갱신됩니다.",
            "여기서 blocker가 남아 있으면 아직 실주문 단계로 가지 않습니다.",
        ],
        0.75,
        1.7,
        4.8,
        2.2,
    )
    _add_image(slide, bootstrap_image, 5.25, 1.45, 7.05)
    _add_badge(slide, "추천 명령", 0.78, 4.28, 1.5, (20, 184, 166))
    _add_bullets(
        slide,
        [
            r".\bootstrap_small_live_validation.cmd -ConfigPath config.live.micro.json -Market KRW-BTC",
        ],
        0.78,
        4.8,
        4.25,
        0.9,
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    _add_title(slide, "4. readiness 확인 후 소액 live", "이 단계부터는 실제 돈이 움직일 수 있으니 아주 작게만")
    _add_bullets(
        slide,
        [
            "터미널 A는 `listen-private`로 체결과 자산 변화를 봅니다.",
            "터미널 B는 `run-live-daemon`으로 아주 짧게 돌립니다.",
            "실제 주문은 1건만, 아주 작은 금액만, 시장/수량/리스크를 눈으로 다시 확인한 뒤 진행합니다.",
        ],
        0.75,
        1.7,
        4.8,
        2.2,
    )
    _add_image(slide, live_run_image, 5.25, 1.45, 7.05)
    _add_bullets(
        slide,
        [
            "체결되면 OK",
            "미체결이면 `open-orders -> order-show -> cancel-order` 순서로 취소 확인",
            "끝나면 `live-reconcile` 한 번 더 실행",
        ],
        0.78,
        4.55,
        4.25,
        1.45,
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    _add_title(slide, "5. 마지막 증빙 저장", "이 3개가 있어야 체크리스트를 닫을 수 있습니다")
    _add_bullets(
        slide,
        [
            "live session report",
            "support bundle zip",
            "release pack status 또는 release verify 상태",
        ],
        0.75,
        1.7,
        3.5,
        1.6,
    )
    _add_image(slide, evidence_image, 4.2, 1.45, 8.05)
    _add_bullets(
        slide,
        [
            "다 끝나면 `PRODUCT_COMPLETION_CHECKLIST.md`의 `Small live validation completed` 아래에 증빙 경로를 적습니다.",
            "그때 제가 체크리스트도 바로 체크해드리면 됩니다.",
        ],
        0.78,
        4.95,
        4.0,
        1.1,
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    _add_title(slide, "현재 상태 요약", "지금 readiness 파일 기준으로는 아직 아래 blocker가 남아 있습니다")
    _add_image(slide, blocker_image, 0.78, 1.35, 6.0)
    _add_bullets(
        slide,
        [
            "남은 blocker가 사라질 때까지는 `bootstrap_small_live_validation.cmd`와 `.env`, `config.live.micro.json`만 먼저 만지면 됩니다.",
            "지금 기준 다음 순서:",
        ]
        + ["- {0}".format(item) for item in next_steps[:5]],
        6.95,
        1.55,
        5.0,
        4.6,
    )

    prs.save(output_path)

    manifest = {
        "pptx_path": str(output_path),
        "asset_dir": str(ASSET_ROOT),
        "readiness_path": str(readiness_path),
    }
    with open(DIST_ROOT / "guide-manifest.json", "w", encoding="utf-8") as handle:
        json.dump(manifest, handle, ensure_ascii=False, indent=2)

    return manifest


def main() -> int:
    parser = argparse.ArgumentParser(description="Generate an easy PPT guide for small live validation.")
    parser.add_argument(
        "--readiness",
        default="dist/live-validation/small-live-validation-readiness.json",
        help="Readiness JSON generated by prepare_small_live_validation.",
    )
    parser.add_argument(
        "--output",
        default="dist/small-live-validation-guide/Upbit-Small-Live-Validation-Guide.pptx",
        help="Output PPTX path.",
    )
    args = parser.parse_args()

    readiness_path = PROJECT_ROOT / args.readiness if not Path(args.readiness).is_absolute() else Path(args.readiness)
    output_path = PROJECT_ROOT / args.output if not Path(args.output).is_absolute() else Path(args.output)

    if not readiness_path.exists():
        raise SystemExit("readiness file not found: {0}".format(readiness_path))

    manifest = build_presentation(output_path=output_path, readiness_path=readiness_path)
    print(json.dumps(manifest, ensure_ascii=False, indent=2))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
